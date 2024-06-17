import os
import openai
from pptx import Presentation

# Load your OpenAI API key
openai.api_key = 'your-openai-api-key'

# Load the files in the results folder
results_folder = "path-to-your-results-folder"
files = os.listdir(results_folder)

# Initialize a Presentation object
prs = Presentation()

# Iterate through each file
for file in files:

    # Open the file and read its content
    with open(os.path.join(results_folder, file), 'r') as f:
        content = f.read()
        
    # you can split the content into chunks if it's too long 
    # content_parts = split_into_chunks(content)

    # Then for each chunk generate a slide
    # for part in content_parts:

    # Generate the presentation text using OpenAI API
    response = openai.Completion.create(
        engine="text-davinci-002",
        prompt=content,
        temperature=0.5,
        max_tokens=100
    )
    
    # Create a new slide with a title and content
    slide_layout = prs.slide_layouts[1]  # 0 is title and content
    slide = prs.slides.add_slide(slide_layout)

    # Add the title and the content
    title = slide.shapes.title
    title.text = "Title for this Slide"

    content = slide.placeholders[1]
    content.text = response.choices[0].text.strip()

# Save the presentation
prs.save('generated_presentation.pptx')